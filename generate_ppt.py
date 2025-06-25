from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import pandas as pd

# Create a presentation object
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Environmental Risk Classification: Interpretation"
subtitle = slide.placeholders[1]
subtitle.text = "Automated Output Interpretation | Submission Date: June 25, 2025"

# Slide: App Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "App Overview"
content = slide.placeholders[1].text_frame
content.text = (
    "The Streamlit app visualizes and classifies environmental risk for mining zones using geographic and pollution data. "
    "It provides interactive risk flags, summary tables, and visualizations."
)

# Slide: Data Snapshot
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Data Snapshot"
df = pd.read_csv("zone_features.csv")
sample = df.head(5).to_string(index=False)
content = slide.placeholders[1].text_frame
content.text = "Sample rows from zone_features.csv:\n\n" + sample

# Slide: Feature Engineering
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Feature Engineering"
content = slide.placeholders[1].text_frame
content.text = (
    "- Water level ranges split into min/max\n"
    "- Risk classified as Low (0) or High (1) using threshold 327\n"
    "- Features: Lat, Long, NO₂, Water, Urban Growth, etc.\n"
    "- Target: Risk_Class (binary)"
)

# Slide: Model & Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Model & Results"
content = slide.placeholders[1].text_frame
content.text = (
    "- Logistic Regression with balanced class weights\n"
    "- Test accuracy: 83.3%\n"
    "- Class 0 (Low Risk): perfect recall, some false positives\n"
    "- Class 1 (High Risk): perfect precision, but missed some positives\n"
    "- Model is cautious, suitable for risk-sensitive applications."
)

# Slide: Insights
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Insights"
content = slide.placeholders[1].text_frame
content.text = (
    "- Dataset is imbalanced\n"
    "- Location features may influence risk\n"
    "- Environmental variables show moderate correlation with risk\n"
    "- Model performance is strong for small data size"
)

# Slide: Submission Note
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Submission Note"
content = slide.placeholders[1].text_frame
content.text = (
    "This PPT was auto-generated from code and data.\n"
    "For further details, see app.py and classification.ipynb."
)

prs.save("output_interpretation_slides.pptx")
print("Presentation generated: output_interpretation_slides.pptx")
